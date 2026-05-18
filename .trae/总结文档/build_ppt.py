from pptx import Presentation
from pptx.util import Pt, Inches, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
DARK_BLUE = RGBColor(0x1A, 0x2A, 0x4A)
ACCENT_BLUE = RGBColor(0x2D, 0x6D, 0xCF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF2, 0xF5)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MID_GRAY = RGBColor(0x66, 0x66, 0x66)
TAG_GREEN = RGBColor(0x27, 0xAE, 0x60)
TAG_ORANGE = RGBColor(0xE6, 0x7E, 0x22)
TAG_RED = RGBColor(0xC0, 0x39, 0x2B)

def add_bg(slide, color=DARK_BLUE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=Pt(18), color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_multiline(slide, left, top, width, height, lines):
    """lines = list of (text, font_size, color, bold, alignment)"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        text, font_size, color, bold, *rest = line_data
        align = rest[0] if rest else PP_ALIGN.LEFT
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = font_size
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = 'Microsoft YaHei'
        p.alignment = align
        p.space_after = Pt(6)
    return tf

def add_card(slide, left, top, width, height, title, body_lines, title_color=ACCENT_BLUE):
    """Add a card-style box with title and body"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.fill.background()
    shape.shadow.inherit = False

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.15)

    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.color.rgb = title_color
    p.font.bold = True
    p.font.name = 'Microsoft YaHei'

    for line in body_lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(10)
        p.font.color.rgb = DARK_GRAY
        p.font.name = 'Microsoft YaHei'
        p.space_before = Pt(2)

# ==================== SLIDE 1: Title ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_text_box(slide, 1.5, 1.8, 10, 1.2, '基于 Spring Boot 的智能宿舍管理系统', Pt(40), WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 3.0, 10, 0.8, '设计与实现', Pt(28), RGBColor(0x7F, 0xB3, 0xFF), False, PP_ALIGN.CENTER)
add_multiline(slide, 1.5, 4.5, 10, 1.5, [
    ('答辩人：王和友    指导教师：刘海波    人工智能学院    计算机科学与技术', Pt(16), RGBColor(0xBB, 0xCC, 0xDD), False, PP_ALIGN.CENTER),
    ('2026年5月', Pt(14), RGBColor(0x99, 0xAA, 0xBB), False, PP_ALIGN.CENTER),
])

# ==================== SLIDE 2: TOC ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 1, 0.5, 5, 0.6, '目录 CONTENTS', Pt(28), DARK_BLUE, True, PP_ALIGN.LEFT)
toc_items = [
    ('01', '研究背景与意义', '传统宿舍管理痛点与数字化转型需求'),
    ('02', '系统需求分析', '三类角色用例与核心功能梳理'),
    ('03', '系统总体设计', 'B/S架构、功能模块与数据库设计'),
    ('04', '关键技术栈', 'Spring Boot + Vue 3 + MySQL + AI集成'),
    ('05', '系统功能实现', '学生端、宿管端、管理员端核心功能'),
    ('06', '系统测试', '功能测试、性能测试、安全与兼容性测试'),
    ('07', '总结与展望', '项目成果、创新点与未来改进方向'),
]
for idx, (num, title, desc) in enumerate(toc_items):
    y = 1.6 + idx * 0.75
    add_text_box(slide, 1.5, y, 0.8, 0.5, num, Pt(20), ACCENT_BLUE, True, PP_ALIGN.LEFT)
    add_text_box(slide, 2.5, y, 4, 0.3, title, Pt(16), DARK_GRAY, True, PP_ALIGN.LEFT)
    add_text_box(slide, 2.5, y + 0.3, 8, 0.3, desc, Pt(10), MID_GRAY, False, PP_ALIGN.LEFT)

# ==================== SLIDE 3: Research Background ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 1, 0.5, 5, 0.6, '01  研究背景与意义', Pt(28), DARK_BLUE, True, PP_ALIGN.LEFT)
add_card(slide, 0.8, 1.5, 3.6, 2.2, '🔴 传统管理困境', [
    '• 高校扩招，学生数量激增',
    '• 宿舍分配、水电抄表依赖人工',
    '• 信息更新滞后，数据统计困难',
    '• 报修线下登记，进度难追踪',
])
add_card(slide, 4.8, 1.5, 3.6, 2.2, '🟡 学生体验痛点', [
    '• 查询住宿信息需去管理处',
    '• 报修无反馈渠道',
    '• 水电费缴纳流程繁琐',
    '• 换寝、访客预约靠线下沟通',
])
add_card(slide, 8.8, 1.5, 3.6, 2.2, '🟢 数字化转型必然', [
    '• 信息技术推动校园管理升级',
    '• 前后端分离降低开发维护成本',
    '• AI + 数据可视化助力科学决策',
    '• 实现宿舍管理的规范化与高效化',
])
add_text_box(slide, 1, 4.3, 11, 0.4, '▎设计目标：将传统人工管理模式中的宿舍分配、报修处理、费用统计等业务流程全面数字化', Pt(13), MID_GRAY, False)

# ==================== SLIDE 4: Use Case Analysis ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 1, 0.5, 5, 0.6, '02  系统需求分析 — 用例图', Pt(28), DARK_BLUE, True, PP_ALIGN.LEFT)
add_card(slide, 0.8, 1.5, 3.6, 3.0, '学生端', [
    '• 登录与个人信息管理',
    '• 在线报修（提交/追踪/评价/取消）',
    '• 水电费查看与线上缴费',
    '• 晚归打卡与补卡申请',
    '• 换寝申请与进度查询',
    '• 访客预约、失物招领',
    '• 智能客服（AI + 人工）',
    '• 宿舍公约、文明宿舍查看',
], TAG_GREEN)
add_card(slide, 4.8, 1.5, 3.6, 3.0, '宿管员端', [
    '• 学生信息管理（增删改查）',
    '• 床位分配与退寝办理',
    '• 报修工单接单、转派、完成',
    '• 卫生检查打分（1-10 分）',
    '• 访客预约审批',
    '• 水电用量录入与催缴',
    '• 打卡统计与补卡审批',
    '• 文明宿舍评选',
], TAG_ORANGE)
add_card(slide, 8.8, 1.5, 3.6, 3.0, '管理员端', [
    '• 三类用户全生命周期管理',
    '• 宿舍楼/房间/床位资源管理',
    '• 报修类型管理与全局查看',
    '• 公告发布、水电阈值配置',
    '• 系统数据统计（ECharts）',
    '• Excel 批量导入导出',
    '• 操作日志与登录日志查看',
    '• 客服聊天管理',
], TAG_RED)
add_text_box(slide, 1, 5.0, 11, 0.4, '▎共 29 张数据库表，覆盖三类角色全业务流程 | SQL 脚本：dormitory_system.sql', Pt(11), MID_GRAY, False)

# ==================== SLIDE 5: System Architecture ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 1, 0.5, 5, 0.6, '03  系统总体设计 — 架构', Pt(28), DARK_BLUE, True, PP_ALIGN.LEFT)
add_card(slide, 0.8, 1.5, 3.6, 2.0, '前端 Frontend', [
    'Vue 3 + Element Plus + Vite',
    '端口 5173 → 代理 → 后端 8080',
    'Axios 统一封装请求',
    'Pinia 管理全局状态',
    '路由懒加载，首屏优化',
])
add_card(slide, 4.8, 1.5, 3.6, 2.0, '后端 Backend (8080)', [
    'Spring Boot 3.2 + JDK 17',
    'Controller → Mapper 简化架构',
    'Spring Security + JWT 认证',
    'MyBatis-Plus 分页 + 批量查询',
    'HikariCP 连接池（最大 50）',
])
add_card(slide, 8.8, 1.5, 3.6, 2.0, '数据库 MySQL 8.0', [
    '29 张表，无物理外键',
    '逻辑关联 via ID',
    'BCrypt 密码加密存储',
    '唯一索引：学号、用户名',
    '支持事务、备份恢复',
])
add_text_box(slide, 1, 4.0, 11, 0.4, '▎B/S 前后端分离架构 | 启动类：DormitoryApplication.java | 配置：application.yml', Pt(11), MID_GRAY, False)

# Modules overview
add_card(slide, 0.8, 4.6, 3.6, 2.3, '学生端（15 页）', [
    '报修 / 水电费 / 打卡 / 换寝',
    '智能客服 / 失物招领',
    '访客预约 / 紧急求助',
    '文明宿舍 / 宿舍公约 / 公告',
], TAG_GREEN)
add_card(slide, 4.8, 4.6, 3.6, 2.3, '宿管端（12 页）', [
    '学生管理 / 床位分配',
    '报修处理 / 卫生检查',
    '水电录入 / 打卡统计',
    '换寝审批 / 访客审批',
], TAG_ORANGE)
add_card(slide, 8.8, 4.6, 3.6, 2.3, '管理员端（25 页）', [
    '用户管理 / 宿舍资源管理',
    '报修类型 / 公告 / 水电阈值',
    '数据统计 / Excel 导入导出',
    '客服管理 / 日志查看',
], TAG_RED)

# ==================== SLIDE 6: Database Design ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 1, 0.5, 5, 0.6, '03  数据库设计', Pt(28), DARK_BLUE, True, PP_ALIGN.LEFT)
add_card(slide, 0.8, 1.5, 5.5, 2.5, '核心实体关系', [
    '用户类：admin、dormitory_manager、student',
    '资源类：building（楼）→ room（房间）→ bed（床位）',
    '学生入住床位，building_id/room_id 关联',
    'bed.status：0-空闲 1-占用 2-维修中 3-损坏待修',
    '无物理外键，代码层通过 MyBatis-Plus 维护关联',
], ACCENT_BLUE)
add_card(slide, 6.8, 1.5, 5.5, 2.5, '业务表一览（共 29 张）', [
    '报修：repair / repair_type / repair_comment',
    '水电费：utility_bill / utility_threshold',
    '打卡：check_in / check_in_apply / check_in_setting',
    '换寝：room_change / check_out',
    '客服：chat_session / chat_message',
    '其他：visitor / emergency_help / lost_and_found / health_check / notice / message',
])
add_card(slide, 0.8, 4.4, 5.5, 2.3, '实体类对应', [
    '每个表对应 dormitory-admin/.../entity/ 下一个 Java 类',
    '如 Student.java → student 表，字段一一映射',
    'MyBatis-Plus BaseMapper 提供默认 CRUD',
    '复杂查询用 LambdaQueryWrapper 构建条件',
], ACCENT_BLUE)
add_card(slide, 6.8, 4.4, 5.5, 2.3, '索引策略', [
    'student.student_number：唯一索引（学号唯一）',
    'building_id / room_id 等关联字段：普通索引',
    'repair.status / repair.student_id：组合查询优化',
    'SQL 建表脚本：项目根目录 dormitory_system.sql',
], ACCENT_BLUE)

# ==================== SLIDE 7: Tech Stack ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 1, 0.5, 5, 0.6, '04  关键技术栈', Pt(28), DARK_BLUE, True, PP_ALIGN.LEFT)
add_card(slide, 0.8, 1.5, 3.6, 2.5, '后端 Spring Boot 3.2', [
    '• 约定优于配置，自动装配',
    '• Spring Security + JWT 认证',
    '• MyBatis-Plus 简化 CRUD',
    '• HikariCP 连接池',
    '• 内嵌 Tomcat，开箱即用',
], ACCENT_BLUE)
add_card(slide, 4.8, 1.5, 3.6, 2.5, '前端 Vue 3 + Element Plus', [
    '• Composition API',
    '• Pinia 状态管理',
    '• Axios 请求统一封装',
    '• ECharts 数据可视化',
    '• 路由懒加载，按需加载',
], TAG_GREEN)
add_card(slide, 8.8, 1.5, 3.6, 2.5, '数据库 MySQL 8.0', [
    '• 开源关系型数据库',
    '• 29 张业务表',
    '• 支持事务、索引优化',
    '• Navicat 可视化管理',
], TAG_ORANGE)
add_card(slide, 0.8, 4.3, 3.6, 2.5, '安全认证体系', [
    '• BCrypt 密码加密（salt + 慢哈希）',
    '• JWT 无状态 Token 认证',
    '• SecurityConfig.java 过滤器链',
    '• 数据层细粒度权限控制',
], TAG_RED)
add_card(slide, 4.8, 4.3, 3.6, 2.5, '创新：DeepSeek AI 集成', [
    '• Java HttpClient 调用 API',
    '• chat_session + chat_message 存储',
    '• 智能客服 + 人工客服双模式',
    '• ChatController.java 实现',
], ACCENT_BLUE)
add_card(slide, 8.8, 4.3, 3.6, 2.5, '工具与效率', [
    '• Apache POI：Excel 导入导出',
    '• JMeter：性能并发测试',
    '• Lombok：简化实体类代码',
    '• TRAE：AI 辅助开发工具',
], MID_GRAY)

# ==================== SLIDE 8: Student Features ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 1, 0.5, 8, 0.6, '05  系统功能实现 — 学生端', Pt(28), DARK_BLUE, True, PP_ALIGN.LEFT)
add_card(slide, 0.8, 1.5, 5.8, 2.5, '在线报修', [
    '• 状态机模式：0-待处理 → 1-已接单 → 2-维修中 → 3-已完成',
    '• 支持故障类型选择、图片上传（最多 4 张）',
    '• 紧急工单置顶标红，优先处理',
    '• 宿管未接单前学生可取消',
    '• 完成后学生评价（满意/一般/不满意）',
    '📄 RepairController.java / RepairCommentController.java',
], TAG_GREEN)
add_card(slide, 7.0, 1.5, 5.5, 2.5, '换寝申请', [
    '• 选择目标宿舍和床位 → 提交申请',
    '• 宿管审批通过后多表联动更新：',
    '  ①释放原床位 ②占用新床位 ③更新 student 表',
    '  ④宿舍人数 ±1 ⑤申请状态 → 已通过',
    '• 所有操作用数据库事务保证一致性',
    '📄 RoomChangeController.java',
], TAG_ORANGE)
add_card(slide, 0.8, 4.3, 5.8, 2.5, '智能客服（创新功能）', [
    '• 对接 DeepSeek API：https://api.deepseek.com/v1/chat/completions',
    '• 智能客服（chatType=1）：AI 自动回复',
    '• 人工客服（chatType=2）：转管理员处理',
    '• 支持多轮对话、会话管理',
    '• 超时 30s，AI 不可用时提示转人工',
    '📄 ChatController.java',
], ACCENT_BLUE)
add_card(slide, 7.0, 4.3, 5.5, 2.5, '其他学生端功能', [
    '• 打卡：22:00-23:00 打卡，超时补卡申请',
    '• 水电费：查看账单，扫码模拟支付',
    '• 访客预约：填写访客信息，宿管审批',
    '• 紧急求助：一键推送宿管端标红提醒',
    '• 失物招领、宿舍公约、文明宿舍查看',
    '📄 CheckInController / UtilityController / VisitorController',
], MID_GRAY)

# ==================== SLIDE 9: Manager & Admin Features ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 1, 0.5, 8, 0.6, '05  系统功能实现 — 宿管端 & 管理员端', Pt(28), DARK_BLUE, True, PP_ALIGN.LEFT)
add_card(slide, 0.8, 1.5, 5.8, 2.5, '宿管端核心功能', [
    '• 学生管理：增删改查本楼栋学生，导出 Excel',
    '• 床位管理：分配床位 → 床位状态更新，支持批量分配',
    '• 报修处理：接单、转派、完成标记，工单状态同步',
    '• 卫生检查：1-10 分打分，拍照留证，支持历史查询',
    '• 水电管理：录入本月用量 → 系统自动计算费用',
    '📄 ManagerController / HealthCheckController / BedController',
], TAG_ORANGE)
add_card(slide, 7.0, 1.5, 5.5, 2.5, '宿管端创新功能', [
    '• 水电费按宿舍类型区别定价（4 人间 / 6 人间不同单价）',
    '• 超限自动预警：用量超阈值 → 创建预警 → 通知学生',
    '• 文明宿舍自动排名：卫生检查均分 → 降序排名',
    '• Excel 批量导入水电用量',
    '📄 UtilityController / UtilityThresholdController',
], TAG_ORANGE)
add_card(slide, 0.8, 4.3, 5.8, 2.5, '管理员端核心功能', [
    '• 用户管理：三类用户 CRUD，启用/禁用，批量导入',
    '• 宿舍资源：楼栋/房间/床位全局管理，床位状态联动',
    '• 报修管理：类型配置、全局查看、数据统计（柱状图）',
    '• 系统设置：公告管理、操作日志、数据备份',
    '📄 AdminController / BuildingController / StatisticsController',
], TAG_RED)
add_card(slide, 7.0, 4.3, 5.5, 2.5, '管理员端数据能力', [
    '• 数据统计：学生总数、入住率、报修分布（ECharts）',
    '• 水电统计：各楼栋用量柱状图、缴费/欠费扇形图',
    '• Excel 导出：学生、报修、打卡数据一键导出',
    '• Excel 导入：学生批量导入、水电费批量导入',
    '• 所有导入导出基于 Apache POI（XSSFWorkbook）',
    '📄 StatisticsController.java',
], TAG_RED)

# ==================== SLIDE 10: N+1 Optimization ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 1, 0.5, 8, 0.6, '05  性能优化 — N+1 查询问题', Pt(28), DARK_BLUE, True, PP_ALIGN.LEFT)
add_card(slide, 0.8, 1.5, 5.8, 2.8, '❌ 优化前（N+1 问题）', [
    '1. 分页查学生（1 次 SQL）',
    '2. 全表查所有床位 selectList(null)（1 次）',
    '3. 循环为每个床位查房间（~1000 次）',
    '4. 循环为每个房间查楼栋（~500 次）',
    '总查询：~1501 次 → 页面加载 2 秒',
    '',
    '📄 StudentController.java list() 方法（旧版）',
], TAG_RED)
add_card(slide, 7.0, 1.5, 5.5, 2.8, '✅ 优化后（批量查询 + 内存关联）', [
    '1. 分页查学生（1 次）',
    '2. 只查当前页学生的床位，用 IN 条件（1 次）',
    '3. 批量查房间 selectBatchIds(roomIds)（1 次）',
    '4. 批量查楼栋 selectBatchIds(buildingIds)（1 次）',
    '总查询：4 次 → 响应 ~50ms',
    '',
    '📄 StudentController.java 第 188-236 行',
], TAG_GREEN)
add_text_box(slide, 1, 4.6, 11, 0.4, '▎核心思路：按需查询（只查当前页） + 批量查询（IN / selectBatchIds） + 内存关联（Map<Long, Room>）', Pt(13), MID_GRAY, False)
add_card(slide, 0.8, 5.2, 5.8, 1.5, '优化效果', [
    '查询次数：1501 次 → 4 次',
    '响应时间：2000ms → ~50ms',
    '用户体验：明显卡顿 → 流畅',
], TAG_GREEN)
add_card(slide, 7.0, 5.2, 5.5, 1.5, '相关文档', [
    '优化方案文档：项目根目录/学生列表查询性能优化方案.md',
    '注意：list() 方法已优化，部分辅助方法待后续优化',
], MID_GRAY)

# ==================== SLIDE 11: Testing ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 1, 0.5, 5, 0.6, '06  系统测试', Pt(28), DARK_BLUE, True, PP_ALIGN.LEFT)
add_card(slide, 0.8, 1.5, 3.0, 2.5, '🟢 功能测试（黑盒）', [
    '• 9 个核心用例',
    '• 报修提交/审批/评价全流程',
    '• 换寝申请/审批/更新全流程',
    '• 客服 AI+人工双模式',
    '• 结果：全部通过',
], TAG_GREEN)
add_card(slide, 4.2, 1.5, 3.0, 2.5, '🔵 性能测试（JMeter）', [
    '• 登录 1000 并发：平均 105ms',
    '• 查询 500 并发：平均 14ms',
    '• 数据提交 1000 并发：平均 1.36s',
    '• 连接池 50，CSV 参数化',
    '• 零错误，零数据丢失',
], ACCENT_BLUE)
add_card(slide, 7.6, 1.5, 3.0, 2.5, '🟡 安全性测试', [
    '• 密码 BCrypt 加密存储 ✓',
    '• JWT Token 防篡改验证 ✓',
    '• 未登录访问 → 401/403 ✓',
    '• SQL 注入：MyBatis-Plus 预编译 ✓',
    '• XSS：Vue 3 默认转义 ✓',
], TAG_ORANGE)
add_card(slide, 11.0, 1.5, 1.5, 2.5, '🟣 兼容性', [
    '• Chrome ✓',
    '• Edge ✓',
    '• Firefox ✓',
    '• 多分辨率',
], MID_GRAY)
add_text_box(slide, 1, 4.5, 11, 0.4, '▎测试环境：Windows 11  |  MySQL 8.0  |  JDK 17  |  JMeter 5.6.3  |  TRAE 开发工具', Pt(11), MID_GRAY, False)

# ==================== SLIDE 12: Summary & Innovation ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 1, 0.5, 5, 0.6, '07  总结与展望', Pt(28), DARK_BLUE, True, PP_ALIGN.LEFT)
add_card(slide, 0.8, 1.5, 3.8, 2.5, '项目成果', [
    '• 29 张表完整数据库设计',
    '• 28 个后端 Controller',
    '• 52 个前端页面（3 角色）',
    '• 报修/换寝/水电/打卡等全流程',
    '• 功能测试全部通过，性能达标',
], TAG_GREEN)
add_card(slide, 5.0, 1.5, 3.8, 2.5, '核心创新点', [
    '• DeepSeek AI 智能客服',
    '• 水电费按宿舍类型定价 + 超限预警',
    '• N+1 查询优化（1501 次 → 4 次）',
    '• Excel 批量导入导出',
    '• 文明宿舍自动评分排名',
], ACCENT_BLUE)
add_card(slide, 9.2, 1.5, 3.6, 2.5, '技术亮点', [
    '• BCrypt + JWT 安全认证',
    '• 换寝多表联动事务保证',
    '• 数据层细粒度权限控制',
    '• 路由懒加载前端优化',
    '• HikariCP 连接池优化',
], TAG_ORANGE)
add_card(slide, 0.8, 4.3, 5.5, 2.5, '未来改进方向', [
    '• 移动端适配：开发微信小程序或 App',
    '• 缓存优化：引入 Redis 热点数据缓存',
    '• 安全增强：部署 HTTPS + @PreAuthorize 注解',
    '• 消息推送：WebSocket 实时通知替代轮询',
    '• 物联网集成：智能水电表、门禁接入',
    '• 大数据分析：住宿行为分析与决策支持',
], MID_GRAY)
add_card(slide, 6.8, 4.3, 5.5, 2.5, '项目收获', [
    '• 完整体验了前后端分离项目的全流程开发',
    '• 掌握了 Spring Boot + Vue 3 实战技术',
    '• 实践了数据库设计、性能优化、安全认证',
    '• 理解了三类角色业务流的设计与实现',
    '• 通过 AI 辅助编程提升了开发效率',
], MID_GRAY)

# ==================== SLIDE 13: Thank You ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_text_box(slide, 1.5, 1.5, 10, 1, '致 谢', Pt(48), WHITE, True, PP_ALIGN.CENTER)
add_multiline(slide, 1.5, 3.0, 10, 3, [
    ('感谢指导教师刘海波老师在课题研究和论文撰写过程中的悉心指导', Pt(16), RGBColor(0xBB, 0xCC, 0xDD), False, PP_ALIGN.CENTER),
    ('', Pt(8), WHITE, False, PP_ALIGN.CENTER),
    ('感谢人工智能学院各位老师的培养与帮助', Pt(16), RGBColor(0xBB, 0xCC, 0xDD), False, PP_ALIGN.CENTER),
    ('', Pt(8), WHITE, False, PP_ALIGN.CENTER),
    ('感谢同学们在项目开发和大学生活中的支持与陪伴', Pt(16), RGBColor(0xBB, 0xCC, 0xDD), False, PP_ALIGN.CENTER),
    ('', Pt(16), WHITE, False, PP_ALIGN.CENTER),
    ('恳请各位老师批评指正！', Pt(20), RGBColor(0x7F, 0xB3, 0xFF), True, PP_ALIGN.CENTER),
])

# Save
output = r"C:\Users\32010\Desktop\AI项目\TRAE\.trae\答辩PPT_精简版.pptx"
prs.save(output)
print(f"Done: {output}")
print(f"Total slides: {len(prs.slides)}")
