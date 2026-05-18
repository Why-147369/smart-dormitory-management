import copy
from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.dml.color import RGBColor
import os

src_path = r"C:\Users\32010\Desktop\AI项目\TRAE\.trae\基于Spring Boot的智能宿舍管理系统设计与实现.pptx"
prs = Presentation(src_path)

def find_text_in_slide(slide, keyword):
    """Find all shapes containing keyword"""
    results = []
    for j, shape in enumerate(slide.shapes):
        if shape.has_text_frame:
            text = shape.text_frame.text
            if keyword in text:
                results.append((j, shape, text[:100]))
    return results

def replace_text_in_shape(shape, old, new):
    """Replace text in a shape's text frame"""
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            if old in para.text:
                # Replace in runs
                for run in para.runs:
                    if old in run.text:
                        run.text = run.text.replace(old, new)
                # Also handle the paragraph-level text
                if old in para.text:
                    # Already handled via runs
                    pass

def replace_text_all_shapes(slide, old, new):
    """Replace text in all shapes of a slide"""
    for shape in slide.shapes:
        if shape.has_text_frame:
            tf = shape.text_frame
            for para in tf.paragraphs:
                for run in para.runs:
                    if old in run.text:
                        run.text = run.text.replace(old, new)

# ============================================================
# SLIDE-BY-SLIDE MODIFICATIONS
# ============================================================

# --- SLIDE 5 (index 4): 可行性分析 → change "IDEA、VS Code" to "TRAE" ---
slide5 = prs.slides[4]
replace_text_all_shapes(slide5, 'IDEA、VS Code', 'TRAE（AI 辅助编程工具）')
replace_text_all_shapes(slide5, 'IDEA & VS Code', 'TRAE')
print("Slide 5: Fixed tool names")

# --- SLIDE 11 (index 10): 关键技术栈-2 → add code file hints ---
slide11 = prs.slides[10]
print(f"Slide 11 shapes: {len(slide11.shapes)}")
for j, shape in enumerate(slide11.shapes):
    if shape.has_text_frame:
        txt = shape.text_frame.text[:80]
        print(f"  [{j}] {txt}")

# --- SLIDE 12 (index 11): 学生端功能 → add controller hints ---
slide12 = prs.slides[11]
print(f"\nSlide 12 shapes: {len(slide12.shapes)}")
for j, shape in enumerate(slide12.shapes):
    if shape.has_text_frame:
        txt = shape.text_frame.text[:80]
        print(f"  [{j}] {txt}")

# --- SLIDE 15 (index 14): 系统测试 → fix tools + data ---
slide15 = prs.slides[14]
replace_text_all_shapes(slide15, 'IDEA &、VS Code', '')
replace_text_all_shapes(slide15, 'trae', 'TRAE（开发工具）')
replace_text_all_shapes(slide15, 'IDEA', '')
replace_text_all_shapes(slide15, '防越权', 'Token 防篡改')
# Try to fix performance data
replace_text_all_shapes(slide15, '平均响应时间 < 2秒', '登录 1000并发 平均 105ms，查询 500并发 平均 14ms')
print("Slide 15: Fixed tools, security, performance")

# --- SLIDE 16 (index 15): 总结 → fix 67x claim ---
slide16 = prs.slides[15]
replace_text_all_shapes(slide16, '67倍', '375 倍（查询次数）')
replace_text_all_shapes(slide16, '查询效率提升了67倍', '查询从 1501 次降到 4 次')
print("Slide 16: Fixed N+1 numbers")

# --- Delete slides ---
# We need to delete: Slide 4 (index 3: 研究意义), Slide 5 alternate, etc.
# Let's identify which slides to delete based on content
slides_to_delete = []

for i, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text
            # Slide 4: 研究意义与设计目的 (redundant with slide 3)
            if '研究意义' in text and '设计目的' in text:
                slides_to_delete.append(i)
                print(f"Mark DELETE slide {i+1}: 研究意义与设计目的")
                break
            # Slide 5: 可行性分析
            if '可行性分析' in text and '技术可行性' in text:
                slides_to_delete.append(i)
                print(f"Mark DELETE slide {i+1}: 可行性分析")
                break

# Actually delete in reverse order
for idx in sorted(slides_to_delete, reverse=True):
    rId = prs.slides._sldIdLst[idx].get('id')
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[idx]
    print(f"Deleted slide at index {idx}")

# --- Merge slide 10+11 (技术栈 two pages) ---
# Keep slide 10, copy DeepSeek content from slide 11 to slide 10
slide10 = prs.slides[9]  # 关键技术栈 page 1
slide11_content = prs.slides[10]  # 关键技术栈 page 2

# Find DeepSeek text in slide 11 and add a note to slide 10
for shape in slide11_content.shapes:
    if shape.has_text_frame:
        txt = shape.text_frame.text
        if 'DeepSeek' in txt:
            # Add a small text box to slide10
            txBox = slide10.shapes.add_textbox(Inches(1), Inches(6), Inches(11), Inches(0.8))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = '创新集成：DeepSeek AI 智能客服（ChatController.java）|  MyBatis-Plus 简化 CRUD |  JWT + BCrypt 安全认证'
            p.font.size = Pt(11)
            p.font.name = 'Microsoft YaHei'
            print("Slide 10: Added DeepSeek note from slide 11")
        if 'JWT' in txt or 'Spring Security' in txt:
            print(f"Slide 11 JWT content: {txt[:100]}")

# Delete slide 11 (now merged into slide 10)
rId = prs.slides._sldIdLst[10].get('id')
prs.part.drop_rel(rId)
del prs.slides._sldIdLst[10]
print("Deleted slide 11 (merged into slide 10)")

# --- Merge slide 17+18+19 (总结展望 + 致谢 + QA) into two: 总结展望 + 致谢 ---
# After deletions, slides may have shifted. Let's find current indices.
for i, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.has_text_frame:
            txt = shape.text_frame.text
            if '致 谢' in txt and len(txt) < 10:
                # Found Thank You slide
                # Check if there's a separate QA slide after it
                print(f"Thank you slide at index {i}")
            if 'Q & A' in txt or '感谢聆听' in txt:
                print(f"QA slide at index {i}: {txt[:50]}")

# Save modified PPT
output_path = r"C:\Users\32010\Desktop\AI项目\TRAE\.trae\答辩PPT_优化版.pptx"
prs.save(output_path)
print(f"\nDone! Saved to: {output_path}")
print(f"Final slide count: {len(prs.slides)}")
